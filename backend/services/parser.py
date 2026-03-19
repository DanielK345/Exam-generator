import pdfplumber
from pptx import Presentation


def parse_pdf(file_path: str) -> list[dict]:
    """Extract text from PDF at page level."""
    pages = []
    with pdfplumber.open(file_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if text and text.strip():
                pages.append({
                    "source": f"page_{i + 1}",
                    "content": text.strip()
                })
    return pages


def parse_pptx(file_path: str) -> list[dict]:
    """Extract text from PPTX at slide level."""
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append({
                "source": f"slide_{i + 1}",
                "content": "\n".join(texts)
            })
    return slides


def parse_file(file_path: str, filename: str) -> list[dict]:
    """Route to the correct parser based on file extension."""
    ext = filename.lower().rsplit(".", 1)[-1]
    if ext == "pdf":
        return parse_pdf(file_path)
    elif ext == "pptx":
        return parse_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")
